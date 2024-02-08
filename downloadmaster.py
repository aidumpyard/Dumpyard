import dash
from dash import dcc, html, Input, Output, State
import plotly.graph_objects as go
import pandas as pd
from dash.exceptions import PreventUpdate
from pptx import Presentation
from pptx.util import Inches
import io
from flask import send_file

app = dash.Dash(__name__)

app.layout = html.Div([
    dcc.Dropdown(
        id='page-selection-dropdown',
        options=[
            {'label': 'Dataframe & Waterfall', 'value': 'page1'},
            {'label': 'Two Waterfalls', 'value': 'page2'},
            {'label': 'Scatter Charts', 'value': 'page3'}
        ],
        value='page1'  # Default value
    ),
    html.Div(id='page-content'),
    html.Button("Download", id="download-button"),
    dcc.Dropdown(
        id='download-selection-dropdown',
        options=[
            {'label': 'Dataframe & Waterfall', 'value': 'page1'},
            {'label': 'Two Waterfalls', 'value': 'page2'},
            {'label': 'Scatter Charts', 'value': 'page3'}
        ],
        value=['page1'],  # Default value
        multi=True
    ),
    dcc.Download(id="download-ppt")
])
@app.callback(
    Output('page-content', 'children'),
    Input('page-selection-dropdown', 'value')
)
def display_page(page):
    if page == 'page1':
        # Code to display dataframe and waterfall chart
        return html.Div([/* Components for Page 1 */])
    elif page == 'page2':
        # Code to display two waterfall charts
        return html.Div([/* Components for Page 2 */])
    elif page == 'page3':
        # Code to display scatter charts
        return html.Div([/* Components for Page 3 */])
    else:
        raise PreventUpdate

@app.callback(
    Output("download-ppt", "data"),
    Input("download-button", "n_clicks"),
    State("download-selection-dropdown", "value"),
    prevent_initial_call=True,
)
def generate_ppt(n_clicks, selected_pages):
    if n_clicks is None:
        raise PreventUpdate
    
    # Initialize PowerPoint presentation
    prs = Presentation()
    
    # Based on selected_pages, add slides
    for page in selected_pages:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Here, add your custom logic to populate the slide based on the page type
        # This may include adding charts, tables, etc.
        # Example: if page == 'page1': add_dataframe_and_waterfall(slide)
    
    # Save presentation to a BytesIO object
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return send_file(
        file_stream,
        as_attachment=True,
        download_name="presentation.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


if __name__ == '__main__':
    app.run_server(debug=True)


