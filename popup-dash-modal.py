import dash
from dash import dcc, html, Input, Output, State
import dash_bootstrap_components as dbc
from pptx import Presentation
import io
from flask import send_file

app = dash.Dash(__name__, external_stylesheets=[dbc.themes.BOOTSTRAP])

app.layout = html.Div([
    dcc.Dropdown(
        id='page-selection-dropdown',
        options=[
            {'label': 'Dataframe & Waterfall', 'value': 'page1'},
            {'label': 'Two Waterfalls', 'value': 'page2'},
            {'label': 'Scatter Charts', 'value': 'page3'}
        ],
        value='page1'  # Default value
    ),
    html.Div(id='page-content'),
    html.Button("Download", id="download-button"),
    dbc.Modal([
        dbc.ModalHeader(dbc.ModalTitle("Select Pages to Download")),
        dbc.ModalBody([
            dcc.Dropdown(
                id='download-selection-dropdown',
                options=[
                    {'label': 'Dataframe & Waterfall', 'value': 'page1'},
                    {'label': 'Two Waterfalls', 'value': 'page2'},
                    {'label': 'Scatter Charts', 'value': 'page3'}
                ],
                value=['page1'],  # Default value
                multi=True
            )
        ]),
        dbc.ModalFooter(
            html.Button("Confirm Download", id="confirm-download", n_clicks=0)
        ),
    ], id="download-modal", is_open=False),
    dcc.Download(id="download-ppt")
])

@app.callback(
    Output('page-content', 'children'),
    Input('page-selection-dropdown', 'value')
)
def display_page(page):
    if page == 'page1':
        return html.Div([
            html.P("Dataframe and Waterfall Model for Page 1"),
        ])
    elif page == 'page2':
        return html.Div([
            html.P("Two Waterfall Models for Page 2"),
        ])
    elif page == 'page3':
        return html.Div([
            html.P("Scatter Charts for Page 3"),
        ])
    else:
        return html.Div([html.P("Select a page to view its content.")])

@app.callback(
    Output("download-modal", "is_open"),
    [Input("download-button", "n_clicks"), Input("confirm-download", "n_clicks")],
    [State("download-modal", "is_open")],
)
def toggle_modal(n1, n2, is_open):
    if n1 or n2:
        return not is_open
    return is_open

@app.callback(
    Output("download-ppt", "data"),
    [Input("confirm-download", "n_clicks")],
    [State("download-selection-dropdown", "value")],
    prevent_initial_call=True,
)
def generate_ppt(n_clicks, selected_pages):
    if not n_clicks:
        return dash.no_update

    prs = Presentation()
    for page in selected_pages:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Placeholder for adding custom slide content based on the selection

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)

    return send_file(
        file_stream,
        as_attachment=True,
        download_name="presentation.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == '__main__':
    app.run_server(debug=True)